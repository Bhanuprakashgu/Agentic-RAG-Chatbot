
from abc import ABC, abstractmethod
from typing import List, Dict, Any

class BaseDocumentParser(ABC):
    """Base class for all document parsers"""
    
    @abstractmethod
    def parse(self, file_path: str) -> Dict[str, Any]:
        """
        Parse a document and return structured content
        
        Returns:
            Dict containing:
            - 'text': extracted text content
            - 'metadata': document metadata (title, author, etc.)
            - 'chunks': list of text chunks
        """
        pass
    
    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks"""
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence endings within the last 100 characters
                last_period = text.rfind('.', start, end)
                last_exclamation = text.rfind('!', start, end)
                last_question = text.rfind('?', start, end)
                
                sentence_end = max(last_period, last_exclamation, last_question)
                if sentence_end > start + chunk_size - 100:
                    end = sentence_end + 1
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - overlap
            
        return chunks





# csv_parser.py
import pandas as pd

class CSVParser(BaseDocumentParser):
    def parse(self, file_path: str) -> Dict[str, Any]:
        df = pd.read_csv(file_path)
        text_content = df.to_string()
        
        return {
            "text": text_content,
            "metadata": {"file_name": file_path, "file_type": "csv"},
            "chunks": self.chunk_text(text_content)
        }




# docx_parser.py
from docx import Document

class DOCXParser(BaseDocumentParser):
    def parse(self, file_path: str) -> Dict[str, Any]:
        """Parse DOCX document"""
        try:
            doc = Document(file_path)
            
            # Extract metadata
            metadata = {
                'title': doc.core_properties.title or 'Unknown',
                'author': doc.core_properties.author or 'Unknown',
                'paragraphs': len(doc.paragraphs),
                'file_type': 'DOCX'
            }
            
            # Extract text from all paragraphs
            text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            # Clean up text
            text = text.strip()
            
            # Create chunks
            chunks = self.chunk_text(text, 1000, 200) # Using default chunk size and overlap for now
            
            return {
                'text': text,
                'metadata': metadata,
                'chunks': chunks
            }
            
        except Exception as e:
            raise Exception(f"Error parsing DOCX: {str(e)}")




# pdf_parser.py
import PyPDF2

class PDFParser(BaseDocumentParser):
    def parse(self, file_path: str) -> Dict[str, Any]:
        """Parse PDF document"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Extract metadata
                metadata = {
                    'title': pdf_reader.metadata.title if pdf_reader.metadata and pdf_reader.metadata.title else 'Unknown',
                    'author': pdf_reader.metadata.author if pdf_reader.metadata and pdf_reader.metadata.author else 'Unknown',
                    'pages': len(pdf_reader.pages),
                    'file_type': 'PDF'
                }
                
                # Extract text from all pages
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                
                # Clean up text
                text = text.strip()
                
                # Create chunks
                chunks = self.chunk_text(text, 1000, 200) # Using default chunk size and overlap for now
                
                return {
                    'text': text,
                    'metadata': metadata,
                    'chunks': chunks
                }
                
        except Exception as e:
            raise Exception(f"Error parsing PDF: {str(e)}")




# pptx_parser.py
from pptx import Presentation

class PPTXParser(BaseDocumentParser):
    def parse(self, file_path: str) -> Dict[str, Any]:
        """Parse PPTX document"""
        try:
            prs = Presentation(file_path)
            
            # Extract metadata
            metadata = {
                'title': prs.core_properties.title or 'Unknown',
                'author': prs.core_properties.author or 'Unknown',
                'slides': len(prs.slides),
                'file_type': 'PPTX'
            }
            
            # Extract text from all slides
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"Slide {slide_num}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                text += slide_text + "\n"
            
            # Clean up text
            text = text.strip()
            
            # Create chunks
            chunks = self.chunk_text(text, 1000, 200) # Using default chunk size and overlap for now
            
            return {
                'text': text,
                'metadata': metadata,
                'chunks': chunks
            }
            
        except Exception as e:
            raise Exception(f"Error parsing PPTX: {str(e)}")




# txt_parser.py
import os

class TXTParser(BaseDocumentParser):
    def parse(self, file_path: str) -> Dict[str, Any]:
        """Parse TXT/Markdown document"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            
            # Extract metadata
            file_name = os.path.basename(file_path)
            file_ext = os.path.splitext(file_name)[1].upper()
            
            metadata = {
                'title': file_name,
                'file_type': file_ext[1:] if file_ext else 'TXT',
                'size': len(text),
                'lines': len(text.split('\n'))
            }
            
            # Clean up text
            text = text.strip()
            
            # Create chunks
            chunks = self.chunk_text(text, 1000, 200) # Using default chunk size and overlap for now
            
            return {
                'text': text,
                'metadata': metadata,
                'chunks': chunks
            }
            
        except Exception as e:
            raise Exception(f"Error parsing TXT/MD: {str(e)}")

def get_parser(file_extension: str) -> BaseDocumentParser:
    """Returns the appropriate parser based on file extension."""
    file_extension = file_extension.lower()
    if file_extension == "pdf":
        return PDFParser()
    elif file_extension == "docx":
        return DOCXParser()
    elif file_extension == "csv":
        return CSVParser()
    elif file_extension in ["txt", "md"]:
        return TXTParser()
    elif file_extension == "pptx":
        return PPTXParser()
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")


